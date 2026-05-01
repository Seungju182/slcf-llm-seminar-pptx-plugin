"""
Smoke test: 모든 의미 단위 메서드와 저수준 메서드를 한 번씩 호출해
'/tmp/slcf_smoke.pptx'에 저장. PowerPoint에서 열어 시각 검증.

uv로 실행:
    uv run --with python-pptx --with pyyaml scripts/smoke_test.py
"""
import sys
from pathlib import Path

# scripts/ 자체를 sys.path에 추가
sys.path.insert(0, str(Path(__file__).parent))

from builder import SeminarBuilder
from style import STYLE


OUTPUT = '/tmp/slcf_smoke.pptx'


def main():
    print(f"STYLE version={STYLE.version}, font={STYLE.fonts.korean}, body={STYLE.sizes.body}")
    print(f"  content area: ({STYLE.grid.content_left():.2f}, {STYLE.grid.content_top():.2f}) "
          f"~ ({STYLE.grid.content_right():.2f}, {STYLE.grid.content_bottom():.2f})")

    b = SeminarBuilder(OUTPUT)

    # 1. 표지
    b.set_cover(date='26.4.30', title='Diffusion Models 정리')

    # 2. 목차 (current=2 → '방법론' 강조)
    b.add_toc(['배경', '방법론', '실험', '결론'], current=2)

    # 3. 챕터 표지
    b.add_section_header(num=1, title='배경', subtitle='Diffusion Models의 등장과 발전')

    # 4. 정의 (definition + why + examples + source_page)
    b.add_definition(
        term='Diffusion Model',
        definition='데이터에 점진적으로 노이즈를 더하고, 그 역과정을 학습해 새 샘플을 생성하는 모델',
        why=[
            'GAN의 학습 불안정 문제 회피',
            '샘플 다양성 우수 (mode collapse 없음)',
            '이론적 해석 명확 (variational bound)',
        ],
        examples=[
            'DDPM (Ho et al., NeurIPS 2020)',
            'DDIM (Song et al., ICLR 2021)',
            'LDM / Stable Diffusion (Rombach et al., CVPR 2022)',
        ],
        source_page=3,
    )

    # 5. 비교표
    b.add_comparison(
        title='주요 모델 비교',
        headers=['모델', 'FID ↓', '샘플링 시간', '특징'],
        rows=[
            ['DDPM', '3.17', '1000 step', '오리지널, 느림'],
            ['DDIM', '4.04', '50 step',   '결정적 샘플링'],
            ['LDM',  '3.60', '50 step',   'latent space 학습'],
            ['EDM',  '1.79', '36 step',   'SOTA-quality'],
        ],
        source_page=12,
    )

    # 6. 절차 (5 steps, '학습' 강조)
    b.add_process(
        title='학습 파이프라인',
        steps=['데이터 수집', '전처리', '학습', '샘플링', '평가'],
        emphasize=[2],
        descriptions=[
            '이미지 큐레이션',
            'normalize, augment',
            'noise schedule + reverse process',
            'denoise from noise',
            'FID, IS, CLIP score',
        ],
        source_page=8,
    )

    # 7. 좌우 두 칸
    b.add_two_content(
        title='장단점',
        left_title='장점',
        left_bullets=[
            '샘플 품질 우수 (FID 1점대)',
            'mode collapse 없음',
            '확률밀도 추정 가능',
            '학습 안정적',
        ],
        right_title='단점',
        right_bullets=[
            '샘플링이 느림 (수십~수백 step)',
            '학습 시간이 김',
            '메모리 요구량 큼',
        ],
    )

    # 8. 결론
    b.add_conclusion(
        takeaways=[
            'Diffusion이 GAN의 안정성 문제를 해결하며 SOTA 달성',
            '샘플 품질·다양성 양쪽에서 우수',
            '다양한 응용 영역으로 빠르게 확산',
        ],
        next_steps=[
            '샘플링 가속화 (DDIM, DPM-Solver, Consistency Models)',
            '조건부 생성 (Classifier-Free Guidance)',
            '비전 외 도메인 (오디오, 비디오, 3D, 분자)',
        ],
    )

    # 9. 기존 저수준 메서드도 정상 동작하는지
    b.add_title_content(
        title='참고: 단순 불릿 본문 슬라이드',
        bullets=[
            '의미 단위 메서드가 안 맞을 때 fallback',
            '제목 + 빨간 불릿 본문',
            'add_title_only로 자유 배치도 가능',
        ],
    )

    # 10. 자유 배치 (add_title_only + 저수준 메서드)
    g = b.style.grid
    s = b.style.sizes
    c = b.style.colors
    slide = b.add_title_only('자유 배치 검증')
    b.add_box(
        slide, '왼쪽 박스 (light_gray)',
        left=g.content_left(), top=g.content_top() + 0.5,
        width=g.content_width() / 2 - 0.2, height=2,
        fill_color=c.light_gray, font_size=s.body,
    )
    b.add_box(
        slide, '오른쪽 박스 (accent_red)',
        left=g.content_left() + g.content_width() / 2 + 0.2,
        top=g.content_top() + 0.5,
        width=g.content_width() / 2 - 0.2, height=2,
        fill_color=c.accent_red, font_size=s.body, bold=True,
    )
    b.add_textbox(
        slide, 'b.style.grid.col(n, of=12) 사용 예시 — 12분할 컬럼 위치 확인',
        left=g.content_left(), top=g.content_top() + 3.0,
        width=g.content_width(), height=0.5,
        font_size=s.body_sub, color=c.muted_gray,
    )
    # 12분할 그리드에 작은 박스 12개 가로 배치
    for i in range(12):
        left, width = g.col(i, of=12)
        b.add_box(
            slide, str(i),
            left=left, top=g.content_top() + 3.7,
            width=width, height=0.5,
            font_size=s.diagram_label,
        )

    out = b.save()
    print(f"\nDone: {out}")

    # 결과 검증
    p = Path(out)
    if not p.exists():
        raise RuntimeError(f"Output not found: {out}")
    size_kb = p.stat().st_size / 1024
    print(f"  size: {size_kb:.1f} KB")

    # 슬라이드 수 체크
    from pptx import Presentation
    prs = Presentation(out)
    print(f"  slides: {len(prs.slides)}")
    for i, sl in enumerate(prs.slides):
        text_shapes = [sh for sh in sl.shapes if sh.has_text_frame]
        first_text = ''
        for sh in text_shapes:
            t = sh.text_frame.text.strip()
            if t:
                first_text = t.split('\n')[0][:50]
                break
        print(f"    slide {i}: {first_text}")


if __name__ == '__main__':
    main()
