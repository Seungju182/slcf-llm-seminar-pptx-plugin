"""
SLCF 세미나 발표 자료 빌더

양식의 슬라이드를 in-place로 deep copy 해서 새 슬라이드를 만드는 방식.
하나의 Presentation 객체만 사용하므로 패키지 충돌 없이 안정적.

사용 예:
    from builder import SeminarBuilder

    b = SeminarBuilder('output.pptx')
    b.set_cover(date='25.5.10', title='Diffusion Models')
    b.add_title_content(
        title='Background',
        bullets=['DDPM (NeurIPS 20)', 'DDIM (ICLR 21)'],
    )
    b.save()

작동 원리:
    1. 양식을 열면 표지(slide 0)와 본문 예시(slide 1)가 들어있음
    2. 빌드 시작 시 이 두 슬라이드를 "원본"으로 기억해두고 hidden 처리
    3. 사용자가 add_*를 호출하면, 해당 원본 슬라이드를 deep copy 해서 새로 추가
    4. save 직전에 원본 두 슬라이드를 제거
"""
import copy
import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree

# style.py와 같은 디렉터리. sys.path 보장을 위해 명시적으로 추가.
sys.path.insert(0, str(Path(__file__).parent))
from style import STYLE, load_style  # noqa: E402

# ===== 호환용 모듈 상수 =====
# 실제 값은 assets/style.yaml에서 로드. 이 상수들은 기존 코드 호환용 alias이며,
# 새로 추가하는 코드에서는 STYLE.colors.body_text 같은 형태로 직접 참조하세요.
# style.yaml을 바꾸면 이 상수들도 자동으로 새 값을 가리킵니다 (재import 시).
BODY_TEXT_COLOR = STYLE.colors.body_text
ACCENT_RED      = STYLE.colors.accent_red
HEADER_BG_BLACK = STYLE.colors.header_bg_black
HEADER_TEXT     = STYLE.colors.header_text
MUTED_GRAY      = STYLE.colors.muted_gray
LIGHT_GRAY      = STYLE.colors.light_gray

KOREAN_FONT = STYLE.fonts.korean
LATIN_FONT  = STYLE.fonts.latin

TEMPLATE_PATH = Path(__file__).parent.parent / 'assets' / 'template.pptx'

# 양식 안에서 어떤 슬라이드가 어떤 역할을 하는지
SOURCE_SLIDE_INDEX = {
    'cover': 0,          # 표지 (어두운 배경 + 뇌 이미지)
    'title_content': 1,  # 본문 (검정 헤더 + 빨간 불릿)
}


def _force_run_color(run_element, rgb=BODY_TEXT_COLOR):
    """run XML에 색상을 강제로 박아넣음. 양식의 기본 파란색을 덮어씀.

    rPr 안에 <a:solidFill><a:srgbClr val="000000"/></a:solidFill> 추가.
    이미 색상이 있으면 교체.
    """
    a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    rPr = run_element.find(f'{{{a_ns}}}rPr')
    if rPr is None:
        rPr = etree.SubElement(run_element, f'{{{a_ns}}}rPr')
        run_element.insert(0, rPr)

    # 기존 solidFill 제거 (있으면)
    for old_fill in rPr.findall(f'{{{a_ns}}}solidFill'):
        rPr.remove(old_fill)
    # 다른 fill 타입도 제거 (gradFill, pattFill 등)
    for fill_tag in ('gradFill', 'pattFill', 'noFill', 'blipFill'):
        for old in rPr.findall(f'{{{a_ns}}}{fill_tag}'):
            rPr.remove(old)

    # 새 solidFill 추가
    solidFill = etree.SubElement(rPr, f'{{{a_ns}}}solidFill')
    srgbClr = etree.SubElement(solidFill, f'{{{a_ns}}}srgbClr')
    srgbClr.set('val', f'{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}')


def _set_text_in_shape(shape, text, force_black=True, color=None):
    """shape의 첫 paragraph 첫 run의 텍스트를 바꿔치기 (서식 유지).

    color (RGBColor, 선택): 명시적 색상. force_black보다 우선.
        예) 강조 슬라이드의 노란 제목.
    force_black=True (기본): 글자색을 검정으로 덮어씀. 본문에 사용.
    force_black=False, color=None: 양식 색상 그대로 (일반 헤더 흰 글씨).
    """
    if not shape.has_text_frame:
        return False
    tf = shape.text_frame
    if not tf.paragraphs:
        tf.text = text
        return True
    p = tf.paragraphs[0]
    if p.runs:
        p.runs[0].text = text
        if color is not None:
            _force_run_color(p.runs[0]._r, color)
        elif force_black:
            _force_run_color(p.runs[0]._r, BODY_TEXT_COLOR)
        for run in p.runs[1:]:
            run.text = ''
    else:
        p.text = text
        if p.runs:
            if color is not None:
                _force_run_color(p.runs[0]._r, color)
            elif force_black:
                _force_run_color(p.runs[0]._r, BODY_TEXT_COLOR)
    return True


def _set_bullets_keep_format(shape, bullets, force_black=True):
    """shape의 첫 paragraph 서식을 복제해서 여러 bullet을 채움.

    force_black=True (기본): 모든 bullet의 글자색을 검정으로 덮어씀.
    """
    if not shape.has_text_frame or not bullets:
        return False
    tf = shape.text_frame
    txBody = tf._txBody

    # 기존 paragraph 모두 제거
    for p in list(txBody.findall(qn('a:p'))):
        txBody.remove(p)

    tf.text = bullets[0]
    first_p = tf.paragraphs[0]
    # bullet 마커/폰트를 paragraph 레벨에서 명시적으로 강제 (slide layout default 차단)
    _force_paragraph_bullet_style(first_p)
    if force_black and first_p.runs:
        _force_run_color(first_p.runs[0]._r, BODY_TEXT_COLOR)

    for bullet in bullets[1:]:
        new_p_element = copy.deepcopy(first_p._p)
        runs = new_p_element.findall(qn('a:r'))
        if runs:
            t = runs[0].find(qn('a:t'))
            if t is not None:
                t.text = bullet
            if force_black:
                _force_run_color(runs[0], BODY_TEXT_COLOR)
            for extra_run in runs[1:]:
                t2 = extra_run.find(qn('a:t'))
                if t2 is not None:
                    t2.text = ''
        txBody.append(new_p_element)

    return True


def _apply_korean_font(run, size_pt=None, bold=False, color=BODY_TEXT_COLOR):
    """run에 한국어 폰트와 색상을 적용. 새로 만든 도형/표의 텍스트에 사용.

    한국어/영문 폰트를 함께 지정해서, 어떤 글자가 와도 깨지지 않게 함.
    """
    run.font.name = LATIN_FONT
    if size_pt is not None:
        run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color

    # 한국어용 ea (East Asian) 폰트도 같이 지정
    a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    rPr = run._r.find(f'{{{a_ns}}}rPr')
    if rPr is None:
        rPr = etree.SubElement(run._r, f'{{{a_ns}}}rPr')
        run._r.insert(0, rPr)
    # 기존 ea 폰트 제거
    for ea in rPr.findall(f'{{{a_ns}}}ea'):
        rPr.remove(ea)
    ea = etree.SubElement(rPr, f'{{{a_ns}}}ea')
    ea.set('typeface', KOREAN_FONT)


def _force_paragraph_bullet_style(paragraph):
    """paragraph의 bullet 마커를 '•' 로, 폰트를 한국어/Latin 통일로 강제.

    `tf.text = "..."` 직후 호출. python-pptx가 paragraph 레벨 pPr을 리셋해
    slide layout의 default bullet (template의 ■ + 모노스페이스)이 노출되는 것을 차단.
    색상/크기/굵기는 건드리지 않음 — _force_run_color 등으로 따로 제어.
    """
    p = paragraph._p
    a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    pPr = p.find(f'{{{a_ns}}}pPr')
    if pPr is None:
        pPr = etree.SubElement(p, f'{{{a_ns}}}pPr')
        p.remove(pPr)
        p.insert(0, pPr)  # pPr은 a:p의 첫 자식이어야 함
    # 기존 bullet 정의 모두 제거
    for tag in ('buChar', 'buAutoNum', 'buNone', 'buBlip', 'buFont'):
        for elem in pPr.findall(f'{{{a_ns}}}{tag}'):
            pPr.remove(elem)
    # buFont (bullet 글자 자체의 폰트) + buChar (마커 문자)
    bu_font = etree.SubElement(pPr, f'{{{a_ns}}}buFont')
    bu_font.set('typeface', LATIN_FONT)
    bu_char = etree.SubElement(pPr, f'{{{a_ns}}}buChar')
    bu_char.set('char', '•')
    # 각 run에 한국어/Latin 폰트 강제 (색/크기/굵기는 보존)
    for run in paragraph.runs:
        run.font.name = LATIN_FONT
        rPr = run._r.find(f'{{{a_ns}}}rPr')
        if rPr is None:
            rPr = etree.SubElement(run._r, f'{{{a_ns}}}rPr')
            run._r.insert(0, rPr)
        for ea in rPr.findall(f'{{{a_ns}}}ea'):
            rPr.remove(ea)
        ea = etree.SubElement(rPr, f'{{{a_ns}}}ea')
        ea.set('typeface', KOREAN_FONT)


class SeminarBuilder:
    """SLCF 세미나 양식을 기반으로 발표 자료를 빌드."""

    def __init__(self, output_path, style_path=None):
        """
        Args:
            output_path: 출력할 .pptx 파일 경로
            style_path: 사용자 정의 style.yaml 경로 (None이면 기본 assets/style.yaml)
        """
        if not TEMPLATE_PATH.exists():
            raise FileNotFoundError(f"양식 파일 없음: {TEMPLATE_PATH}")

        # 인스턴스별 스타일 객체. 기본은 모듈 임포트 시 로드된 STYLE 싱글턴.
        # style_path가 주어지면 새로 로드해서 이 인스턴스에만 적용.
        self.style = load_style(style_path) if style_path else STYLE

        self.prs = Presentation(str(TEMPLATE_PATH))
        self.output_path = output_path

        # 양식의 원본 슬라이드를 "복제 소스"로 기억
        # _slides_to_remove에 모아뒀다가 save 시 제거
        self._source_cover = self.prs.slides[SOURCE_SLIDE_INDEX['cover']]
        self._source_content = self.prs.slides[SOURCE_SLIDE_INDEX['title_content']]
        self._original_slide_ids = [s.slide_id for s in self.prs.slides]

    def _clone_slide(self, source_slide):
        """양식 안 슬라이드를 deep copy 해서 새 슬라이드로 추가.

        모든 shape, 배경, 이미지 관계를 그대로 복제함.
        주의: shape XML 안에 r:embed="rIdN" 같은 참조가 있으면, 새 슬라이드의
        rels에 같은 target_part로 관계를 만들고 그 rId로 교체해야 그림이 보임.
        """
        from pptx.oxml.ns import nsmap

        # 같은 레이아웃으로 새 빈 슬라이드 생성
        new_slide = self.prs.slides.add_slide(source_slide.slide_layout)

        # 새 슬라이드의 기본 placeholder 모두 제거
        for shape in list(new_slide.shapes):
            sp = shape._element
            sp.getparent().remove(sp)

        # 원본 슬라이드의 rId → target_part 맵 만들기
        source_rels = {rel.rId: rel for rel in source_slide.part.rels.values()}

        # 원본의 모든 shape를 복제해서 추가하면서 rId 참조 업데이트
        for shape in source_slide.shapes:
            new_element = copy.deepcopy(shape._element)

            # shape XML 안의 모든 r:embed, r:link 속성 찾아서 rId 재매핑
            r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
            for elem in new_element.iter():
                for attr_name in (f'{{{r_ns}}}embed', f'{{{r_ns}}}link'):
                    old_rid = elem.get(attr_name)
                    if old_rid and old_rid in source_rels:
                        rel = source_rels[old_rid]
                        if "image" in rel.reltype or "chart" in rel.reltype:
                            # 새 슬라이드 rels에 같은 target으로 관계 추가
                            new_rid = new_slide.part.relate_to(
                                rel.target_part, rel.reltype
                            )
                            elem.set(attr_name, new_rid)

            new_slide.shapes._spTree.insert_element_before(new_element, 'p:extLst')

        return new_slide

    def _remove_original_slides(self):
        """양식의 원본 예시 슬라이드들 제거."""
        sldIdLst = self.prs.slides._sldIdLst
        original_ids = set(self._original_slide_ids)
        for sldId in list(sldIdLst):
            # slide id 매칭으로 원본 슬라이드 찾기
            try:
                sid = int(sldId.get('id'))
                if sid in original_ids:
                    rId = sldId.get(qn('r:id'))
                    self.prs.part.drop_rel(rId)
                    sldIdLst.remove(sldId)
            except (TypeError, ValueError):
                pass

    # ===== 공개 메서드 =====

    def set_cover(self, date, title=None):
        """표지 슬라이드 추가.

        Args:
            date: 'YY.M.D' 형식 (예: '25.5.10')
            title: 발표 제목 (선택). 날짜 옆에 추가됨.
        """
        slide = self._clone_slide(self._source_cover)

        for shape in slide.shapes:
            if shape.has_text_frame and '세미나' in shape.text_frame.text:
                # 표지의 날짜 텍스트는 양식 기본이 흰색(어두운 배경 위).
                # force_black=False로 양식 색상을 유지해야 어두운 배경에 묻히지 않음.
                if title:
                    _set_text_in_shape(shape, f"{date} 세미나 — {title}",
                                       force_black=False)
                else:
                    _set_text_in_shape(shape, f"{date} 세미나",
                                       force_black=False)
                break

        return slide

    def add_title_content(self, title, bullets, important=False):
        """제목 + 불릿 본문 슬라이드 (가장 자주 쓰는 형식).

        Args:
            title: 헤더에 들어갈 제목 (예: 'Background')
            bullets: 본문 불릿 리스트
            important: True면 헤더 글씨를 노란색으로 강조
        """
        slide = self._clone_slide(self._source_content)
        title_color = self.style.colors.important_yellow if important else None

        title_set = False
        bullets_set = False
        body_shape = None

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            current = shape.text_frame.text

            if not title_set and current.strip() and '\n' not in current.strip():
                # 일반 헤더는 양식 흰색, important=True이면 노란색으로 덮어씀
                _set_text_in_shape(shape, title, force_black=False,
                                   color=title_color)
                title_set = True
            elif not bullets_set and '\n' in current:
                _set_bullets_keep_format(shape, bullets)
                body_shape = shape
                bullets_set = True

        # 본문 영역을 콘텐츠 그리드 전체로 확장 + bullet 수에 따라 폰트 동적 조정.
        # 양식의 body placeholder(H=2.2")가 너무 작아 슬라이드 위쪽에 콘텐츠가
        # 압축되던 문제를 fix.
        if body_shape is not None:
            g = self.style.grid
            body_shape.left = Inches(g.content_left())
            body_shape.top = Inches(g.header_bottom + 0.25)
            body_shape.width = Inches(g.content_width())
            body_shape.height = Inches(
                g.footer_top - (g.header_bottom + 0.25) - 0.2
            )

            # bullet 수 기반 폰트 (시각적 무게감)
            n = max(1, len(bullets))
            if n <= 2:
                body_size = 28
            elif n == 3:
                body_size = 26
            elif n == 4:
                body_size = 22
            elif n == 5:
                body_size = 20
            else:
                body_size = 18

            tf = body_shape.text_frame
            tf.word_wrap = True
            try:
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE  # 수직 중앙 정렬
            except Exception:
                pass
            for para in tf.paragraphs:
                para.space_before = Pt(8)
                para.space_after = Pt(8)
                for run in para.runs:
                    run.font.size = Pt(body_size)

        # 양식 본문에 있던 그림(image) 제거
        for shape in list(slide.shapes):
            if shape.shape_type == 13:  # PICTURE
                sp = shape._element
                sp.getparent().remove(sp)

        return slide

    def add_title_only(self, title, important=False):
        """제목만 있는 슬라이드 (이미지/도표 직접 배치용).

        Args:
            title: 헤더 제목
            important: True면 헤더 글씨를 노란색으로 강조
        """
        slide = self._clone_slide(self._source_content)
        title_color = self.style.colors.important_yellow if important else None

        title_set = False
        for shape in list(slide.shapes):
            if not shape.has_text_frame:
                if shape.shape_type == 13:  # 그림 제거
                    sp = shape._element
                    sp.getparent().remove(sp)
                continue
            current = shape.text_frame.text
            if not title_set and current.strip() and '\n' not in current.strip():
                _set_text_in_shape(shape, title, force_black=False,
                                   color=title_color)
                title_set = True
            elif '\n' in current:
                # 본문 영역 통째로 제거
                sp = shape._element
                sp.getparent().remove(sp)

        return slide

    def add_image_grid(self, title, images, captions=None, source_page=None,
                       important=False):
        """N개의 이미지를 자동 그리드로 배치. n=1→1칸, 2→1×2, 3→1×3, 4→2×2.

        각 이미지는 cell 안에서 aspect ratio 보존 + 가로 중앙 정렬.
        captions가 있으면 각 이미지 아래 작은 회색 캡션.

        Args:
            title: 슬라이드 제목
            images: 이미지 경로 리스트 (1~4개)
            captions: 각 이미지 캡션 (선택, len(images)와 동일 길이)
            source_page: 원문 페이지 (선택)
            important: True면 헤더 노란색
        """
        n = len(images)
        if n < 1 or n > 4:
            raise ValueError(f"image_grid는 1~4개 이미지만 지원 (받음: {n})")

        slide = self.add_title_only(title, important=important)
        g = self.style.grid
        s = self.style.sizes
        c = self.style.colors

        has_caption = bool(captions and any(captions))
        cap_h = 0.4 if has_caption else 0.0

        area_left = g.content_left()
        area_top = g.content_top() + 0.2
        area_w = g.content_width()
        area_h = g.content_bottom() - area_top - (0.4 if source_page else 0.1)

        if n == 1:
            rows, cols = 1, 1
        elif n == 2:
            rows, cols = 1, 2
        elif n == 3:
            rows, cols = 1, 3
        else:  # n == 4
            rows, cols = 2, 2

        gap = 0.2
        cell_w = (area_w - gap * (cols - 1)) / cols
        cell_h = (area_h - gap * (rows - 1)) / rows
        img_max_h = cell_h - (cap_h + 0.05 if has_caption else 0)

        EMU_PER_INCH = 914400
        for i, img_path in enumerate(images):
            r, col = divmod(i, cols)
            cell_left = area_left + col * (cell_w + gap)
            cell_top = area_top + r * (cell_h + gap)

            # 자연 크기로 추가 → cell에 맞춰 aspect-preserve 리사이즈
            pic = slide.shapes.add_picture(
                img_path, Inches(cell_left), Inches(cell_top)
            )
            ratio = pic.height / pic.width if pic.width else 1.0
            target_w_emu = Inches(cell_w)
            target_h_emu = int(target_w_emu * ratio)
            max_h_emu = Inches(img_max_h)
            if target_h_emu > max_h_emu:
                target_h_emu = max_h_emu
                target_w_emu = int(target_h_emu / ratio) if ratio else max_h_emu
            pic.width = target_w_emu
            pic.height = target_h_emu
            # cell 안에서 가로 중앙 정렬 (top은 cell 상단)
            actual_w_inch = target_w_emu / EMU_PER_INCH
            pic.left = Inches(cell_left + (cell_w - actual_w_inch) / 2)
            pic.top = Inches(cell_top)

            if has_caption and i < len(captions) and captions[i]:
                actual_h_inch = target_h_emu / EMU_PER_INCH
                cap_top = cell_top + actual_h_inch + 0.05
                self.add_textbox(
                    slide, captions[i],
                    left=cell_left, top=cap_top,
                    width=cell_w, height=cap_h,
                    font_size=s.caption, color=c.muted_gray, align='center'
                )

        if source_page:
            self._add_source_caption(slide, source_page)

        return slide

    def add_image(self, slide, image_path, left, top, width=None, height=None):
        """슬라이드에 이미지 추가. 위치/크기는 inch.

        가급적 외부 이미지(PDF 캡처 등)에 의존하지 말고, add_table/add_box로
        직접 그리세요. 이미지가 꼭 필요할 때만 (사진, 복잡한 도식 등) 사용.
        """
        kwargs = {'left': Inches(left), 'top': Inches(top)}
        if width is not None:
            kwargs['width'] = Inches(width)
        if height is not None:
            kwargs['height'] = Inches(height)
        return slide.shapes.add_picture(image_path, **kwargs)

    def add_table(self, slide, data, left, top, width, height,
                  header_row=True, header_col=False, font_size=None,
                  header_font_size=None):
        """양식 색상에 맞는 표 추가 (검정 헤더 + 한국어 폰트).

        Args:
            slide: 슬라이드 객체
            data: 2차원 리스트. 예: [['항목', '값'], ['정확도', '95%']]
                  모든 텍스트는 한국어로 작성 권장.
            left, top, width, height: 위치/크기 (inch 단위)
            header_row: 첫 행을 헤더 스타일(검정 배경 + 흰 글씨)로
            header_col: 첫 열을 헤더 스타일로
            font_size: 본문 폰트 크기 (pt). None이면 STYLE.sizes.table_body
            header_font_size: 헤더 폰트 크기 (pt). None이면 STYLE.sizes.table_header

        Returns:
            추가된 table 객체
        """
        if font_size is None:
            font_size = self.style.sizes.table_body
        if header_font_size is None:
            header_font_size = self.style.sizes.table_header

        rows = len(data)
        cols = len(data[0]) if rows > 0 else 0
        if rows == 0 or cols == 0:
            raise ValueError("data는 비어있지 않은 2차원 리스트여야 합니다")

        table_shape = slide.shapes.add_table(
            rows, cols,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        table = table_shape.table

        for r, row_data in enumerate(data):
            for c, cell_text in enumerate(row_data):
                cell = table.cell(r, c)
                is_header = (header_row and r == 0) or (header_col and c == 0)

                # 셀 배경
                if is_header:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = HEADER_BG_BLACK
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

                # 텍스트
                tf = cell.text_frame
                tf.text = str(cell_text)
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                if p.runs:
                    text_color = HEADER_TEXT if is_header else BODY_TEXT_COLOR
                    cell_font_size = header_font_size if is_header else font_size
                    _apply_korean_font(
                        p.runs[0],
                        size_pt=cell_font_size,
                        bold=is_header,
                        color=text_color
                    )
                # 셀 vertical center
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        return table

    def add_box(self, slide, text, left, top, width, height,
                fill_color=None, text_color=None, font_size=None, bold=False,
                align='center'):
        """양식 색상에 맞는 박스(사각형) + 텍스트 추가.

        프로세스 다이어그램, 강조 영역, 단계 표시 등에 사용.
        한국어 폰트가 자동 적용됨.

        Args:
            slide: 슬라이드 객체
            text: 박스 안 텍스트 (한국어 권장)
            left, top, width, height: 위치/크기 (inch)
            fill_color: 배경색 RGBColor. None이면 흰색 + 검정 테두리
            text_color: 글자색. None이면 검정 (배경이 어두우면 자동으로 흰색)
            font_size: 폰트 크기 (pt). None이면 STYLE.sizes.body_sub
            bold: 굵게 여부
            align: 'left', 'center', 'right'

        Returns:
            추가된 shape 객체
        """
        if font_size is None:
            font_size = self.style.sizes.body_sub
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )

        # 배경색
        if fill_color is None:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            shape.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
            auto_text_color = BODY_TEXT_COLOR
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
            shape.line.fill.background()
            # 배경이 어두우면 글자는 자동으로 흰색
            brightness = (fill_color[0] * 299 + fill_color[1] * 587 + fill_color[2] * 114) / 1000
            auto_text_color = HEADER_TEXT if brightness < 128 else BODY_TEXT_COLOR

        if text_color is None:
            text_color = auto_text_color

        # 텍스트
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.text = text
        p = tf.paragraphs[0]
        align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
        p.alignment = align_map.get(align, PP_ALIGN.CENTER)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        if p.runs:
            _apply_korean_font(
                p.runs[0], size_pt=font_size, bold=bold, color=text_color
            )

        return shape

    def add_textbox(self, slide, text, left, top, width, height,
                    font_size=None, bold=False, color=None, align='left'):
        """투명 배경의 텍스트박스 추가. 자유 위치 텍스트 배치용.

        한국어 폰트가 모든 줄에 자동 적용됨 (텍스트에 \\n이 있어도 일관).
        font_size None이면 STYLE.sizes.body_sub (14pt 기본).
        """
        if font_size is None:
            font_size = self.style.sizes.body_sub
        from pptx.util import Inches
        tb = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        tf.text = text
        align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER,
                     'right': PP_ALIGN.RIGHT}
        # \n으로 만들어진 모든 paragraph에 동일 서식 적용 (첫 줄만 적용되는 버그 방지)
        for p in tf.paragraphs:
            p.alignment = align_map.get(align, PP_ALIGN.LEFT)
            for run in p.runs:
                _apply_korean_font(
                    run,
                    size_pt=font_size,
                    bold=bold,
                    color=color or BODY_TEXT_COLOR
                )
        return tb

    # ===== 의미 단위 슬라이드 메서드 =====
    # 좌표·폰트·색이 메서드 안에 잠겨있어, 호출자는 콘텐츠만 신경 씀.
    # 모든 위치는 self.style.grid에서 계산. 폰트는 self.style.sizes에서.

    def _add_source_caption(self, slide, source_page):
        """슬라이드 우하단에 '(원문 p.N)' 캡션 추가."""
        g = self.style.grid
        s = self.style.sizes
        c = self.style.colors
        return self.add_textbox(
            slide, f'(원문 p.{source_page})',
            left=g.slide_w - 2.0,
            top=g.footer_top - 0.05,
            width=1.4, height=0.3,
            font_size=s.caption, color=c.muted_gray, align='right'
        )

    def add_speaker_notes(self, slide, text):
        """슬라이드의 발표자 노트(presenter notes)에 텍스트 설정.

        PowerPoint의 'Notes' pane에 표시. 발표 시 발표자만 보는 대본.
        기존 노트가 있으면 덮어씀.
        """
        if not text:
            return None
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.text = text
        return notes_tf

    def add_section_header(self, num, title, subtitle='', important=False):
        """챕터/섹션 표지 슬라이드.

        본문 헤더바에 'N. {title}'을 넣고, 가운데에 큰 'Chapter N'(빨강)과
        제목(검정 굵게)을 배치. 발표 흐름의 chapter 전환점.

        Args:
            num: 챕터 번호 (int)
            title: 챕터 제목
            subtitle: 부제 (선택)
            important: True면 헤더 글씨를 노란색으로 강조
        """
        slide = self.add_title_only(f'{num}. {title}', important=important)
        g = self.style.grid
        s = self.style.sizes
        c = self.style.colors

        self.add_textbox(slide, f'Chapter {num}',
            left=g.content_left(), top=2.5,
            width=g.content_width(), height=1.4,
            font_size=s.section_number, bold=True,
            color=c.accent_red, align='center')

        self.add_textbox(slide, title,
            left=g.content_left(), top=4.0,
            width=g.content_width(), height=1.0,
            font_size=s.section_title, bold=True,
            color=c.body_text, align='center')

        if subtitle:
            self.add_textbox(slide, subtitle,
                left=g.content_left(), top=5.2,
                width=g.content_width(), height=0.6,
                font_size=s.body, color=c.muted_gray, align='center')

        return slide

    def add_toc(self, items, current=None, title='목차', important=False):
        """목차 슬라이드. 번호+항목으로 세로 정렬. current 지정 시 해당 행 강조.

        Args:
            items: 목차 항목 리스트 (한국어)
            current: 현재 진행 중인 항목 (1-based). None이면 강조 없음.
            title: 슬라이드 제목 (기본 '목차')
            important: True면 헤더 글씨를 노란색으로 강조
        """
        slide = self.add_title_only(title, important=important)
        g = self.style.grid
        s = self.style.sizes
        c = self.style.colors

        n = len(items)
        item_h = 0.6
        total_h = n * item_h
        start_top = (g.content_top() + g.content_bottom()) / 2 - total_h / 2 + 0.2

        for i, item in enumerate(items):
            is_current = (current is not None and (i + 1 == current))
            text = f'  {i+1:02d}.    {item}'
            color = c.accent_red if is_current else c.body_text
            self.add_textbox(slide, text,
                left=g.content_left() + 1.5,
                top=start_top + i * item_h,
                width=g.content_width() - 3.0, height=item_h,
                font_size=s.body, bold=is_current, color=color, align='left')

        return slide

    def add_definition(self, term, definition, why=None, examples=None,
                       source_page=None, important=False):
        """정의/개념 소개 슬라이드.

        헤더에 용어, 본문 상단 회색 박스에 정의, 그 아래에
        '왜 필요한가?' / '예시' 섹션 배치. 둘 다 주어지면 좌우 2칸,
        하나만 주어지면 풀폭 단일 섹션 (콘텐츠 영역 안에 항상 fit).

        Args:
            term: 용어 (예: '정칙화')
            definition: 한 문장 정의
            why: 필요 이유 리스트
            examples: 구체 예시 리스트
            source_page: 원문 페이지 번호 (있으면 우하단에 캡션)
            important: True면 헤더 글씨를 노란색으로 강조
        """
        slide = self.add_title_only(term, important=important)
        g = self.style.grid
        s = self.style.sizes
        c = self.style.colors

        box_top = g.content_top() + 0.2
        box_h = 1.0   # 정의 박스를 더 컴팩트하게 → 아래 섹션에 더 많은 공간
        self.add_box(slide, definition,
            left=g.content_left(), top=box_top,
            width=g.content_width(), height=box_h,
            fill_color=c.light_gray, font_size=s.body, align='left')

        sec_top = box_top + box_h + 0.35
        # source caption이 차지할 공간(약 0.4")을 빼고 사용 가능한 높이
        avail_h = g.content_bottom() - sec_top - (0.4 if source_page else 0.1)

        def _draw_section(label, items, left, width, top, height):
            self.add_textbox(slide, label,
                left=left, top=top, width=width, height=0.45,
                font_size=s.body, bold=True, color=c.accent_red)
            text = '\n'.join('• ' + i for i in items)
            self.add_textbox(slide, text,
                left=left + 0.3, top=top + 0.55,
                width=width - 0.3, height=height - 0.55,
                font_size=s.body)  # 14pt(body_sub) → 18pt(body)로 격상

        if why and examples:
            # 좌우 2칸 — 12분할 그리드 6/6
            L_left, L_w = g.span(0, 6, of=12)
            R_left, R_w = g.span(6, 6, of=12)
            _draw_section('왜 필요한가?', why, L_left, L_w, sec_top, avail_h)
            _draw_section('예시', examples, R_left, R_w, sec_top, avail_h)
        elif why:
            _draw_section('왜 필요한가?', why,
                g.content_left(), g.content_width(), sec_top, avail_h)
        elif examples:
            _draw_section('예시', examples,
                g.content_left(), g.content_width(), sec_top, avail_h)

        if source_page:
            self._add_source_caption(slide, source_page)

        return slide

    def add_comparison(self, title, headers, rows, source_page=None,
                       important=False):
        """비교 표 슬라이드. add_table을 콘텐츠 영역 정중앙에 자동 배치.

        Args:
            title: 슬라이드 제목 (예: 'L1 vs L2 정칙화')
            headers: 1행에 들어갈 헤더 리스트
            rows: 2행 이후 데이터 (2차원 리스트)
            source_page: 원문 페이지 (선택)
            important: True면 헤더 글씨를 노란색으로 강조
        """
        slide = self.add_title_only(title, important=important)
        g = self.style.grid

        data = [list(headers)] + [list(r) for r in rows]
        n_rows = len(data)
        row_h = 0.55
        table_h = min(g.content_height() * 0.85, row_h * n_rows + 0.3)
        table_top = g.content_top() + (g.content_height() - table_h) / 2

        self.add_table(slide, data,
            left=g.content_left(), top=table_top,
            width=g.content_width(), height=table_h)

        if source_page:
            self._add_source_caption(slide, source_page)
        return slide

    def add_process(self, title, steps, emphasize=None, descriptions=None,
                    source_page=None, important=False):
        """절차/파이프라인 다이어그램 슬라이드. N개 박스를 가로로 균등 배치 + 화살표.

        Args:
            title: 슬라이드 제목
            steps: 단계 텍스트 리스트 (3-7개 권장)
            emphasize: 강조할 단계 인덱스 리스트 (0-based). 빨간 배경.
            descriptions: 각 단계 아래 보충 설명 (parallel list, 빈 문자열 허용)
            source_page: 원문 페이지 (선택)
            important: True면 헤더 글씨를 노란색으로 강조
        """
        slide = self.add_title_only(title, important=important)
        g = self.style.grid
        s = self.style.sizes
        c = self.style.colors
        sp = self.style.spacing
        emphasize = set(emphasize or [])

        n = len(steps)
        if n == 0:
            return slide

        arrow_w = sp.process_arrow_w
        box_h = sp.process_box_h
        total_w = g.content_width()
        box_w = (total_w - (n - 1) * arrow_w) / n
        box_top = g.content_top() + (g.content_height() - box_h) / 2 - 0.4

        for i, step in enumerate(steps):
            x = g.content_left() + i * (box_w + arrow_w)
            if i in emphasize:
                fill = c.accent_red
            elif i == n - 1 and not emphasize:
                fill = c.header_bg_black
            else:
                fill = None
            bold = (i in emphasize) or (i == n - 1)
            self.add_box(slide, step,
                left=x, top=box_top, width=box_w, height=box_h,
                fill_color=fill, font_size=s.body, bold=bold)

            if i < n - 1:
                self.add_textbox(slide, '→',
                    left=x + box_w, top=box_top + box_h * 0.15,
                    width=arrow_w, height=box_h * 0.7,
                    font_size=s.arrow, color=c.muted_gray, align='center')

            if descriptions and i < len(descriptions) and descriptions[i]:
                self.add_textbox(slide, descriptions[i],
                    left=x, top=box_top + box_h + 0.15,
                    width=box_w, height=0.6,
                    font_size=s.diagram_label,
                    color=c.muted_gray, align='center')

        if source_page:
            self._add_source_caption(slide, source_page)
        return slide

    def add_two_content(self, title, left_title, left_bullets,
                        right_title, right_bullets, source_page=None,
                        important=False):
        """좌우 두 칸 본문 슬라이드. 12분할 그리드의 좌측 6칸 / 우측 6칸.

        Args:
            title: 슬라이드 제목
            left_title, left_bullets: 좌측 소제목 + 불릿 리스트
            right_title, right_bullets: 우측 소제목 + 불릿 리스트
            source_page: 원문 페이지 (선택)
            important: True면 헤더 글씨를 노란색으로 강조
        """
        slide = self.add_title_only(title, important=important)
        g = self.style.grid
        s = self.style.sizes
        c = self.style.colors

        L_left, L_w = g.span(0, 6, of=12)
        R_left, R_w = g.span(6, 6, of=12)
        top = g.content_top() + 0.3
        avail_h = g.content_height() - 1.0

        for col_left, col_w, col_title, col_bullets in [
            (L_left, L_w, left_title, left_bullets),
            (R_left, R_w, right_title, right_bullets),
        ]:
            self.add_textbox(slide, col_title,
                left=col_left, top=top, width=col_w, height=0.5,
                font_size=s.body, bold=True, color=c.accent_red, align='left')
            self.add_textbox(slide,
                '\n'.join('• ' + b for b in col_bullets),
                left=col_left, top=top + 0.6,
                width=col_w, height=avail_h,
                font_size=s.body_sub, align='left')

        if source_page:
            self._add_source_caption(slide, source_page)
        return slide

    def add_conclusion(self, takeaways, next_steps=None, title='결론',
                       important=False):
        """결론 슬라이드. '핵심 정리' / '향후 연구' 두 섹션.

        Args:
            takeaways: 핵심 정리 불릿 리스트
            next_steps: 향후 연구/방향 불릿 (선택)
            title: 슬라이드 제목 (기본 '결론')
            important: True면 헤더 글씨를 노란색으로 강조
        """
        slide = self.add_title_only(title, important=important)
        g = self.style.grid
        s = self.style.sizes
        c = self.style.colors

        cur_top = g.content_top() + 0.3

        def _section(label, items, body_size, bullet_pad=0.45):
            nonlocal cur_top
            self.add_textbox(slide, label,
                left=g.content_left(), top=cur_top,
                width=g.content_width(), height=0.5,
                font_size=s.body, bold=True, color=c.accent_red)
            cur_top += 0.55
            text = '\n'.join('• ' + i for i in items)
            row_h = bullet_pad * len(items) + 0.3
            self.add_textbox(slide, text,
                left=g.content_left() + 0.3, top=cur_top,
                width=g.content_width() - 0.3, height=row_h,
                font_size=body_size)
            cur_top += row_h + 0.3

        _section('핵심 정리', takeaways, s.body)
        if next_steps:
            _section('향후 연구', next_steps, s.body)  # 14pt → 18pt 일관성

        return slide

    def save(self):
        """저장 (양식 원본 슬라이드 제거 후)."""
        self._remove_original_slides()
        os.makedirs(os.path.dirname(os.path.abspath(self.output_path)) or '.', exist_ok=True)
        self.prs.save(self.output_path)
        print(f"저장 완료: {self.output_path}")
        return self.output_path


if __name__ == '__main__':
    print("SLCF 세미나 양식 빌더 — 사용 예는 모듈 docstring 참고")
